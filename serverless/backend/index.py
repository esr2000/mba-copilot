"""MBA Copilot - FastAPI Backend.

A RAG-powered document Q&A system for MBA students.
"""

from __future__ import annotations

import csv
import io
import os
import random
import string
import time
from datetime import datetime, timezone
from typing import TYPE_CHECKING, Annotated, Any, cast

import fitz  # PyMuPDF
import tiktoken
from docx import Document
from dotenv import load_dotenv
from fastapi import FastAPI, File, Form, HTTPException, UploadFile
from fastapi.middleware.cors import CORSMiddleware
from pptx import Presentation
from pydantic import BaseModel

if TYPE_CHECKING:
    from openai import OpenAI
    from openai.types.chat import ChatCompletionMessageParam
    from pinecone import Index

# =============================================================================
# App
# =============================================================================
load_dotenv()
app = FastAPI(title="MBA Copilot API", root_path="/backend")

app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],  # tighten in prod
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

# =============================================================================
# Configuration
# =============================================================================


class Config:
    """Configuration settings for the MBA Copilot application."""

    # OpenAI
    OPENAI_API_KEY = os.environ.get("OPENAI_API_KEY")
    OPENAI_BASE_URL = os.environ.get("OPENAI_BASE_URL", "https://api.openai.com/v1")
    EMBEDDING_MODEL = "text-embedding-3-large"
    CHAT_MODEL = "gpt-4o-mini"
    EMBEDDING_DIMENSIONS = 1024

    # Pinecone
    PINECONE_API_KEY = os.environ.get("PINECONE_API_KEY")
    PINECONE_INDEX = os.environ.get("PINECONE_INDEX", "mba-copilot")

    # RAG Settings (token-based)
    CHUNK_TOKENS_DOCS = 800
    CHUNK_OVERLAP_TOKENS_DOCS = 150

    # Retrieval settings
    RETRIEVAL_TOP_K = 20
    CONTEXT_MAX_CHUNKS = 8
    MIN_SCORE = 0.25


config = Config()

# =============================================================================
# Clients (lazy init)
# =============================================================================

_openai_client: OpenAI | None = None
_pinecone_index: Index | None = None


def get_openai() -> OpenAI:
    """Get or initialize the OpenAI client."""
    global _openai_client
    if _openai_client is None:
        if not config.OPENAI_API_KEY:
            raise RuntimeError("OPENAI_API_KEY is not set")

        from openai import OpenAI

        _openai_client = OpenAI(
            api_key=config.OPENAI_API_KEY,
            base_url=config.OPENAI_BASE_URL,
        )

    return _openai_client


def get_pinecone_index() -> Index:
    """Get or initialize the Pinecone index."""
    global _pinecone_index
    if _pinecone_index is None:
        if not config.PINECONE_API_KEY:
            raise RuntimeError("PINECONE_API_KEY is not set")

        from pinecone import Pinecone

        pc = Pinecone(api_key=config.PINECONE_API_KEY)
        _pinecone_index = pc.Index(config.PINECONE_INDEX)

    return _pinecone_index


# =============================================================================
# Token Utilities
# =============================================================================


def num_tokens(text: str, model: str | None = None) -> int:
    """Count tokens in text using tiktoken."""
    encoding_model = model or config.EMBEDDING_MODEL
    try:
        enc = tiktoken.encoding_for_model(encoding_model)
    except KeyError:
        enc = tiktoken.get_encoding("cl100k_base")
    return len(enc.encode(text))


def chunk_by_tokens(
    text: str,
    chunk_tokens: int,
    overlap_tokens: int,
    model: str | None = None,
) -> list[str]:
    """Split text into chunks by token count (not characters)."""
    text = text.replace("\r\n", "\n").strip()
    if not text:
        return []

    if overlap_tokens >= chunk_tokens:
        raise ValueError("overlap_tokens must be < chunk_tokens")

    encoding_model = model or config.EMBEDDING_MODEL
    try:
        enc = tiktoken.encoding_for_model(encoding_model)
    except KeyError:
        enc = tiktoken.get_encoding("cl100k_base")

    tokens = enc.encode(text)

    if len(tokens) <= chunk_tokens:
        return [text]

    chunks: list[str] = []
    start = 0

    while start < len(tokens):
        end = min(start + chunk_tokens, len(tokens))
        chunk_text = enc.decode(tokens[start:end]).strip()
        if chunk_text:
            chunks.append(chunk_text)

        if end >= len(tokens):
            break

        start = max(0, end - overlap_tokens)

    return chunks


# =============================================================================
# Document Processing
# =============================================================================


def _extract_pdf_text_best_fidelity(content: bytes) -> str:
    """Extract text from PDF with best fidelity using block coordinates."""
    doc = fitz.open(stream=content, filetype="pdf")
    try:
        pages_out: list[str] = []
        y_tol = 3.0

        for page in doc:
            blocks: Any = page.get_text("blocks")
            clean_blocks: list[Any] = []

            for b in blocks:
                if (
                    isinstance(b, (tuple, list))
                    and len(b) >= 5
                    and isinstance(b[4], str)
                    and b[4].strip()
                ):
                    clean_blocks.append(b)

            clean_blocks.sort(key=lambda b: (round(float(b[1]) / y_tol), float(b[0])))

            page_text = "\n".join(str(b[4]).rstrip() for b in clean_blocks).strip()
            if page_text:
                pages_out.append(page_text)

        return "\n\n".join(pages_out).strip()
    finally:
        doc.close()


def _extract_docx_text(content: bytes) -> str:
    """Extract text from DOCX file."""
    doc = Document(io.BytesIO(content))
    parts: list[str] = []

    for para in doc.paragraphs:
        t = para.text.strip()
        if t:
            parts.append(t)

    for table in doc.tables:
        for row in table.rows:
            line = "\t".join(cell.text.strip() for cell in row.cells).rstrip()
            if line.strip():
                parts.append(line)

    return "\n".join(parts).strip()


def _extract_pptx_text(content: bytes) -> str:
    """Extract text from PPTX file."""
    prs = Presentation(io.BytesIO(content))
    slides_out: list[str] = []

    for si, slide in enumerate(prs.slides, start=1):
        parts: list[str] = [f"--- Slide {si} ---"]

        for shape in slide.shapes:
            s = cast(Any, shape)

            if hasattr(s, "text_frame") and s.text_frame:
                txt = (s.text_frame.text or "").strip()
                if txt:
                    parts.append(txt)

            if hasattr(s, "table") and s.table:
                for row in s.table.rows:
                    line = "\t".join(cell.text.strip() for cell in row.cells).rstrip()
                    if line.strip():
                        parts.append(line)

        try:
            notes_slide = slide.notes_slide
            if notes_slide and notes_slide.notes_text_frame:
                notes_txt = (notes_slide.notes_text_frame.text or "").strip()
                if notes_txt:
                    parts.append("[Notes]\n" + notes_txt)
        except Exception:
            pass

        slides_out.append("\n".join(parts).strip())

    return "\n\n".join(s for s in slides_out if s).strip()


def _extract_pdf_with_pages(content: bytes) -> list[dict[str, Any]]:
    """Extract PDF with page-level metadata: page_number, text."""
    doc = fitz.open(stream=content, filetype="pdf")
    try:
        pages: list[dict[str, Any]] = []
        y_tol = 3.0

        for page_num, page in enumerate(doc, start=1):
            blocks: Any = page.get_text("blocks")
            clean_blocks: list[Any] = []

            for b in blocks:
                if (
                    isinstance(b, (tuple, list))
                    and len(b) >= 5
                    and isinstance(b[4], str)
                    and b[4].strip()
                ):
                    clean_blocks.append(b)

            clean_blocks.sort(key=lambda b: (round(float(b[1]) / y_tol), float(b[0])))

            page_text = "\n".join(str(b[4]).rstrip() for b in clean_blocks).strip()
            if page_text:
                pages.append({"page_number": page_num, "text": page_text})

        return pages
    finally:
        doc.close()


def extract_structured_chunks(file: UploadFile) -> list[dict[str, Any]]:
    """Extract file into structured chunks with metadata.

    - For PDF: chunks include page_number.
    - For others: simple token-based chunking.
    """
    content = file.file.read()
    try:
        file.file.seek(0)
    except Exception:
        pass

    if not isinstance(file.filename, str) or not file.filename:
        raise ValueError("Uploaded file has no filename")

    filename = file.filename.lower()

    # PPTX
    if filename.endswith(".pptx"):
        text = _extract_pptx_text(content)

    # CSV (as plain text)
    elif filename.endswith(".csv"):
        text = content.decode("utf-8-sig", errors="replace")

    # PDF (per-page, keep page_number)
    elif filename.endswith(".pdf"):
        pages = _extract_pdf_with_pages(content)
        if not pages:
            return []

        structured_chunks: list[dict[str, Any]] = []
        for p in pages:
            page_text = str(p.get("text", "")).strip()
            if not page_text:
                continue

            page_chunks = chunk_by_tokens(
                page_text,
                chunk_tokens=config.CHUNK_TOKENS_DOCS,
                overlap_tokens=config.CHUNK_OVERLAP_TOKENS_DOCS,
            )

            for ch in page_chunks:
                structured_chunks.append(
                    {
                        "text": ch,
                        "chunk_index": len(structured_chunks),
                        "page_number": int(p.get("page_number", 0)),
                    }
                )

        return structured_chunks

    # DOCX
    elif filename.endswith(".docx"):
        text = _extract_docx_text(content)

    # TXT/MD
    elif filename.endswith((".txt", ".md")):
        text = content.decode("utf-8-sig", errors="replace")

    else:
        raise HTTPException(status_code=400, detail=f"Unsupported file type: {filename}")

    if not text.strip():
        return []

    text_chunks = chunk_by_tokens(
        text,
        chunk_tokens=config.CHUNK_TOKENS_DOCS,
        overlap_tokens=config.CHUNK_OVERLAP_TOKENS_DOCS,
    )

    return [{"text": chunk, "chunk_index": i} for i, chunk in enumerate(text_chunks)]


def generate_document_id() -> str:
    """Generate a unique document ID."""
    return f"doc_{int(time.time())}_{''.join(random.choices(string.ascii_lowercase, k=6))}"


# =============================================================================
# Embeddings
# =============================================================================


def generate_embedding(text: str) -> list[float]:
    """Generate embedding for a single text string."""
    client = get_openai()
    response = client.embeddings.create(
        model=config.EMBEDDING_MODEL,
        input=text,
        dimensions=config.EMBEDDING_DIMENSIONS,
    )
    return response.data[0].embedding


async def generate_embeddings_batch(texts: list[str]) -> list[list[float]]:
    """Generate embeddings for multiple texts using parallel individual requests."""
    import asyncio
    from openai import AsyncOpenAI

    async_client = AsyncOpenAI(
        api_key=config.OPENAI_API_KEY,
        base_url=config.OPENAI_BASE_URL if config.OPENAI_BASE_URL else None,
    )

    async def get_single_embedding(text: str) -> list[float]:
        response = await async_client.embeddings.create(
            model=config.EMBEDDING_MODEL,
            input=text,
            dimensions=config.EMBEDDING_DIMENSIONS,
        )
        return response.data[0].embedding

    embeddings = await asyncio.gather(*[get_single_embedding(text) for text in texts])
    return list(embeddings)


# =============================================================================
# Pinecone Operations
# =============================================================================


def store_chunks(chunks: list[dict[str, Any]]) -> None:
    """Store document chunks in Pinecone vector database."""
    index = get_pinecone_index()

    batch_size = 100
    for i in range(0, len(chunks), batch_size):
        batch = chunks[i : i + batch_size]
        vectors = [{"id": c["id"], "values": c["embedding"], "metadata": c["metadata"]} for c in batch]
        index.upsert(vectors=vectors)


def query_similar(
    embedding: list[float],
    top_k: int | None = None,
    document_ids: list[str] | None = None,
) -> list[dict[str, Any]]:
    """Query Pinecone for similar document chunks."""
    index = get_pinecone_index()

    query_filter = None
    if document_ids:
        query_filter = {"document_id": {"$in": document_ids}}

    results = index.query(
        vector=embedding,
        top_k=top_k or config.RETRIEVAL_TOP_K,
        include_metadata=True,
        filter=query_filter,
    )

    return [
        {
            "id": m.id,
            "score": m.score,
            "text": (m.metadata or {}).get("text", ""),
            "filename": (m.metadata or {}).get("filename", ""),
            "document_id": (m.metadata or {}).get("document_id", ""),
            "metadata": m.metadata,
        }
        for m in results.matches
    ]


def delete_document(document_id: str) -> None:
    """Delete all chunks for a document from Pinecone."""
    index = get_pinecone_index()
    index.delete(filter={"document_id": {"$eq": document_id}})


def list_documents() -> list[dict[str, Any]]:
    """Best-effort listing using the 'is_first_chunk' marker."""
    index = get_pinecone_index()

    results = index.query(
        vector=[0.0] * config.EMBEDDING_DIMENSIONS,
        top_k=100,
        include_metadata=True,
        filter={"is_first_chunk": {"$eq": True}},
    )

    documents: list[dict[str, Any]] = []
    for m in results.matches:
        md = m.metadata or {}
        documents.append(
            {
                "id": md.get("document_id"),
                "filename": md.get("filename"),
                "chunks": md.get("total_chunks", 1),
                "uploaded_at": md.get("uploaded_at", ""),
            }
        )

    return documents


# =============================================================================
# RAG Pipeline
# =============================================================================


def generate_answer(
    question: str,
    context: str,
    history: list[dict[str, Any]] | None = None,
    chat_model: str | None = None,
    system_prompt: str | None = None,
) -> str:
    """Generate an answer using OpenAI's chat completion API."""
    client = get_openai()

    prompt = system_prompt or "You are a helpful AI assistant."
    model = chat_model or config.CHAT_MODEL

    messages: list[dict[str, str]] = [{"role": "system", "content": prompt}]

    if context:
        messages.append(
            {
                "role": "system",
                "content": (
                    "Here is relevant information from the student's documents:\n\n"
                    f"{context}\n\n"
                    "Use this to answer the question. Cite sources when appropriate."
                ),
            }
        )
    else:
        messages.append(
            {
                "role": "system",
                "content": (
                    "No relevant documents were found. If needed, let the student know "
                    "they should upload relevant materials, but still try to help with general knowledge."
                ),
            }
        )

    if history:
        for msg in history:
            role = str(msg.get("role", ""))
            content = str(msg.get("content", ""))
            if role in ("user", "assistant", "system"):
                messages.append({"role": role, "content": content})

    messages.append({"role": "user", "content": question})

    response = client.chat.completions.create(
        model=model,
        messages=cast("list[ChatCompletionMessageParam]", messages),
        temperature=0.7,
        max_tokens=1000,
    )
    return response.choices[0].message.content or ""


# =============================================================================
# API Models
# =============================================================================


class ChatSettings(BaseModel):
    """Settings for chat completion and RAG retrieval."""
    chat_model: str = "gpt-4o-mini"
    top_k: int = 15  # backwards compat
    min_score: float = 0.3
    system_prompt: str = "You are a helpful AI assistant."


class ChatRequest(BaseModel):
    """Request model for chat endpoint."""
    message: str
    history: list[dict[str, Any]] | None = None
    settings: ChatSettings | None = None
    document_ids: list[str] | None = None


class ChatResponse(BaseModel):
    """Response model for chat endpoint."""
    answer: str
    sources: list[dict[str, Any]]


# =============================================================================
# API Endpoints
# =============================================================================


@app.post("/chat", response_model=ChatResponse)
async def chat(request: ChatRequest) -> ChatResponse:
    """Chat endpoint - answer questions using RAG retrieval."""
    try:
        settings = request.settings or ChatSettings()

        query_embedding = generate_embedding(request.message)

        similar = query_similar(
            query_embedding,
            top_k=config.RETRIEVAL_TOP_K,
            document_ids=request.document_ids,
        )

        relevant = [c for c in similar if float(c.get("score", 0.0)) >= settings.min_score]

        if relevant:
            context_chunks = relevant[: config.CONTEXT_MAX_CHUNKS]
        elif similar:
            context_chunks = similar[: max(3, config.CONTEXT_MAX_CHUNKS // 2)]
        else:
            context_chunks = []

        def _format_source_label(c: dict[str, Any]) -> str:
            md = c.get("metadata") or {}
            fn = c.get("filename") or "unknown"
            page = md.get("page_number")
            slide = md.get("slide_number")

            if page:
                return f"{fn} (p. {page})"
            if slide:
                return f"{fn} (slide {slide})"
            return fn

        if context_chunks:
            context = "\n\n---\n\n".join(
                [f"[Source: {_format_source_label(c)}]\n{c['text']}" for c in context_chunks]
            )
        else:
            context = ""

        answer = generate_answer(
            request.message,
            context,
            request.history,
            chat_model=settings.chat_model,
            system_prompt=settings.system_prompt,
        )

        sources = [
            {
                "text": c["text"],
                "score": c["score"],
                "filename": c["filename"],
                "document_id": c["document_id"],
                "metadata": c.get("metadata", {}),
            }
            for c in context_chunks
        ]

        return ChatResponse(answer=answer, sources=sources)

    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e)) from e


@app.post("/upload")
async def upload(
    file: Annotated[UploadFile, File()],
    filename: Annotated[str | None, Form()] = None,
) -> dict[str, Any]:
    """Upload and process a document file."""
    try:
        if not config.OPENAI_API_KEY:
            raise HTTPException(
                status_code=500,
                detail="OPENAI_API_KEY not configured. Please set environment variables in Vercel dashboard.",
            )
        if not config.PINECONE_API_KEY:
            raise HTTPException(
                status_code=500,
                detail="PINECONE_API_KEY not configured. Please set environment variables in Vercel dashboard.",
            )

        display_filename = filename or file.filename or "unknown"

        structured_chunks = extract_structured_chunks(file)
        if not structured_chunks:
            raise HTTPException(status_code=400, detail="No content to process")

        existing_docs = list_documents()
        for doc in existing_docs:
            if doc.get("filename") == display_filename and doc.get("id"):
                print(f"Deleting existing document with filename: {display_filename}")
                delete_document(str(doc["id"]))

        chunk_texts = [chunk["text"] for chunk in structured_chunks]
        embeddings = await generate_embeddings_batch(chunk_texts)

        document_id = generate_document_id()
        uploaded_at = datetime.now(timezone.utc).isoformat()

        chunks: list[dict[str, Any]] = []
        for i, (structured_chunk, embedding) in enumerate(
            zip(structured_chunks, embeddings, strict=False)
        ):
            chunks.append(
                {
                    "id": f"{document_id}_chunk_{i}",
                    "embedding": embedding,
                    "metadata": {
                        "text": structured_chunk["text"],
                        "document_id": document_id,
                        "filename": display_filename,
                        "page_number": structured_chunk.get("page_number"),
                        "slide_number": structured_chunk.get("slide_number"),
                        "chunk_index": i,
                        "total_chunks": len(structured_chunks),
                        "uploaded_at": uploaded_at,
                        "is_first_chunk": i == 0,
                    },
                }
            )

        store_chunks(chunks)

        return {
            "success": True,
            "document_id": document_id,
            "filename": display_filename,
            "chunks": len(structured_chunks),
        }

    except HTTPException:
        raise
    except Exception as e:
        import traceback

        error_detail = f"{str(e)}\n\nTraceback:\n{traceback.format_exc()}"
        raise HTTPException(status_code=500, detail=error_detail) from e


@app.get("/documents")
async def get_documents() -> dict[str, Any]:
    """Get list of all uploaded documents."""
    try:
        documents = list_documents()
        return {"documents": documents}
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e)) from e


@app.delete("/documents/{document_id}")
async def remove_document(document_id: str) -> dict[str, bool]:
    """Delete a document and all its chunks."""
    try:
        delete_document(document_id)
        return {"success": True}
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e)) from e


@app.post("/upload-from-url")
async def upload_from_url(request: dict[str, Any]) -> dict[str, Any]:
    """Download a file from a URL and process it."""
    try:
        url = request.get("url")
        filename = request.get("filename")

        if not url or not filename:
            raise HTTPException(status_code=400, detail="Missing url or filename")

        if not config.OPENAI_API_KEY:
            raise HTTPException(
                status_code=500,
                detail="OPENAI_API_KEY not configured. Please set environment variables in Vercel dashboard.",
            )
        if not config.PINECONE_API_KEY:
            raise HTTPException(
                status_code=500,
                detail="PINECONE_API_KEY not configured. Please set environment variables in Vercel dashboard.",
            )

        import httpx

        async with httpx.AsyncClient(timeout=300.0) as client:
            response = await client.get(url)
            response.raise_for_status()
            content = response.content

        class FakeUploadFile:
            def __init__(self, content: bytes, filename: str):
                self.file = io.BytesIO(content)
                self.filename = filename
                self.content_type = "application/octet-stream"

        fake_file = FakeUploadFile(content, filename)

        structured_chunks = extract_structured_chunks(fake_file)  # type: ignore[arg-type]
        if not structured_chunks:
            raise HTTPException(status_code=400, detail="No content to process")

        existing_docs = list_documents()
        for doc in existing_docs:
            if doc.get("filename") == filename and doc.get("id"):
                print(f"Deleting existing document with filename: {filename}")
                delete_document(str(doc["id"]))

        chunk_texts = [chunk["text"] for chunk in structured_chunks]
        embeddings = await generate_embeddings_batch(chunk_texts)

        document_id = generate_document_id()
        uploaded_at = datetime.now(timezone.utc).isoformat()

        chunks: list[dict[str, Any]] = []
        for i, (structured_chunk, embedding) in enumerate(
            zip(structured_chunks, embeddings, strict=False)
        ):
            chunks.append(
                {
                    "id": f"{document_id}_chunk_{i}",
                    "embedding": embedding,
                    "metadata": {
                        "text": structured_chunk["text"],
                        "document_id": document_id,
                        "filename": filename,
                        "page_number": structured_chunk.get("page_number"),
                        "slide_number": structured_chunk.get("slide_number"),
                        "chunk_index": i,
                        "total_chunks": len(structured_chunks),
                        "uploaded_at": uploaded_at,
                        "is_first_chunk": i == 0,
                    },
                }
            )

        store_chunks(chunks)

        return {
            "success": True,
            "document_id": document_id,
            "filename": filename,
            "chunks": len(structured_chunks),
        }

    except HTTPException:
        raise
    except Exception as e:
        import traceback

        error_detail = f"{str(e)}\n\nTraceback:\n{traceback.format_exc()}"
        raise HTTPException(status_code=500, detail=error_detail) from e


@app.get("/health")
async def health() -> dict[str, str]:
    """Health check endpoint."""
    return {"status": "ok"}


if __name__ == "__main__":
    import uvicorn

    uvicorn.run(app, host="0.0.0.0", port=8000)
